import sys
sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── colour palette ──────────────────────────────────────────────────────────
NAVY   = RGBColor(0x0A, 0x22, 0x40)
BLUE   = RGBColor(0x00, 0x6E, 0xC0)
ACCENT = RGBColor(0x00, 0xB4, 0xD8)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY  = RGBColor(0xF0, 0xF4, 0xF8)
DGRAY  = RGBColor(0x33, 0x33, 0x33)
MGRAY  = RGBColor(0x88, 0x88, 0x88)
GREEN  = RGBColor(0x1A, 0x6B, 0x3C)
AMBER  = RGBColor(0xD4, 0x6A, 0x10)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

# ── helpers ─────────────────────────────────────────────────────────────────
def add_slide():
    return prs.slides.add_slide(blank)

def rect(slide, l, t, w, h, fill=None, line_clr=None):
    s = slide.shapes.add_shape(1,
        Inches(l), Inches(t), Inches(w), Inches(h))
    s.line.fill.background()
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line_clr:
        s.line.color.rgb = line_clr; s.line.width = Pt(0.75)
    else:
        s.line.fill.background()
    return s

def tb(slide, text, l, t, w, h,
       fs=14, bold=False, italic=False,
       clr=DGRAY, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(
        Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(fs)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = clr
    return box

def multiline(slide, lines, l, t, w, h, fs=12, clr=DGRAY):
    """lines = list of str; leading '  ' → indent"""
    box = slide.shapes.add_textbox(
        Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for line in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = line
        r.font.size = Pt(fs)
        r.font.color.rgb = clr
        first = False
    return box

def header_bar(slide, title, sub=None):
    rect(slide, 0, 0, 13.33, 1.1, fill=NAVY)
    rect(slide, 0, 1.1, 13.33, 0.06, fill=ACCENT)
    tb(slide, title, 0.3, 0.07, 11.5, 0.65,
       fs=28, bold=True, clr=WHITE)
    if sub:
        tb(slide, sub, 0.3, 0.68, 12.5, 0.38,
           fs=13, italic=True, clr=ACCENT)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ════════════════════════════════════════════════════════════════════════════
s = add_slide()
rect(s, 0, 0, 13.33, 7.5, fill=NAVY)
rect(s, 0, 2.5, 13.33, 0.07, fill=ACCENT)
rect(s, 0, 4.62, 13.33, 0.07, fill=ACCENT)

tb(s, 'Ghost Chatbot', 0.5, 0.55, 12.3, 1.3,
   fs=54, bold=True, clr=WHITE, align=PP_ALIGN.CENTER)
tb(s, 'An Educational Narrative Puzzle Game\nfor Teaching IBM SkillsBuild Chatbot Concepts',
   0.8, 1.75, 11.7, 0.85,
   fs=18, clr=ACCENT, align=PP_ALIGN.CENTER)
tb(s, '3rd Group Meeting  ·  MSc Advanced Computing  ·  King\'s College London  ·  March 2026',
   0.5, 2.65, 12.3, 0.45,
   fs=13, italic=True, clr=WHITE, align=PP_ALIGN.CENTER)

pillars = [
    ('Game-Based Learning', '24 levels / 8 Acts\nIntrinsic puzzle-concept\nintegration'),
    ('Narrative Learning',  'Sci-fi mystery story\nGhost grows as learner\nmasters each concept'),
    ('LLM Tutor (Lily)',    'Ollama / llama3.2\nLocal inference\nAdaptive scaffolding'),
]
for i, (title, body) in enumerate(pillars):
    x = 0.45 + i * 4.28
    rect(s, x, 4.78, 3.95, 2.28, fill=RGBColor(0x0F, 0x33, 0x5E))
    tb(s, title, x+0.12, 4.83, 3.7, 0.5,
       fs=14, bold=True, clr=ACCENT)
    tb(s, body, x+0.12, 5.35, 3.7, 1.55,
       fs=12.5, clr=WHITE)

print('Slide 1 done')

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Updates since 2nd Meeting
# ════════════════════════════════════════════════════════════════════════════
s = add_slide()
rect(s, 0, 0, 13.33, 7.5, fill=LGRAY)
header_bar(s, 'Updates Since 2nd Group Meeting', 'Key changes and progress since March 2026 presentation')

updates = [
    (NAVY,  'Game Design Document (GDD v2) — Completed',
     '24 levels across 8 Acts fully specified. Every puzzle mechanic mapped to an IBM SkillsBuild concept.\n'
     "Ghost's linguistic growth arc defined Act-by-Act: system menus → emoji → words → sentences → natural speech."),
    (BLUE,  'Literature Review — 22 verified papers, 9 theoretical sections',
     'GBL, NCLE, Pedagogical Agents, LLM tutors, Chatbot education reviewed. Integrated theoretical framework built.\n'
     'Research gaps clearly identified and argued (Gap 1–4 in §2.10).'),
    (GREEN, 'Methodology — Full Chapter 3 drafted',
     'DBR + quasi-experimental design. 3 RQs operationalised with hypotheses, analysis plan, and OSF pre-registration plan.\n'
     'Isomorphic items, text-volume control (Flesch-Kincaid), and LLM hyperparameter locking all specified.'),
    (AMBER, 'Tech change: IBM Granite  →  Ollama / llama3.2  ⚠',
     'Reason: Local inference eliminates API costs and LLM latency confound in RQ2 evaluation.\n'
     'No watsonx.ai account dependency. Parameters fully reproducible (temperature = 0.3, system prompt pre-registered).'),
]

y = 1.28
for clr, title, body in updates:
    rect(s, 0.35, y, 12.6, 1.35, fill=WHITE)
    rect(s, 0.35, y, 0.14, 1.35, fill=clr)
    tb(s, title, 0.58, y+0.07, 12.15, 0.38,
       fs=14, bold=True, clr=NAVY)
    tb(s, body, 0.58, y+0.5, 12.15, 0.75,
       fs=12, clr=DGRAY)
    y += 1.44

print('Slide 2 done')

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Game Architecture
# ════════════════════════════════════════════════════════════════════════════
s = add_slide()
rect(s, 0, 0, 13.33, 7.5, fill=LGRAY)
header_bar(s, 'Initial Design — Game Architecture',
           '8 Acts × 3 Levels = 24 levels  |  Full IBM SkillsBuild curriculum coverage')

acts = [
    ('Act 0', 'Chatbot Fundamentals',  'Drag-to-match: components overview; what a chatbot is and is not'),
    ('Act 1', 'Intent',                'Sort utterances into intent categories; build intent taxonomy'),
    ('Act 2', 'Entity',                'Annotate text spans; label entity types and synonyms'),
    ('Act 3', 'Dialog',                'Sequence dialog nodes; connect conditions and responses in flow'),
    ('Act 4', 'Confidence & Fallback', 'Slider calibration: set thresholds; design fallback responses'),
    ('Act 5', 'Testing',               'Predict chatbot responses; identify failure cases; write test cases'),
    ('Act 6', 'Integration',           'Assemble full pipeline: intents + entities + dialog + thresholds'),
    ('Act ★', 'Advanced NLP (opt.)',   'Context variables, disambiguation, webhook integration — post-Act 6'),
]

hdrs = ['Act', 'IBM SkillsBuild Topic', 'Puzzle Mechanic (Intrinsic Integration)']
col_x = [0.35, 1.42, 4.12]
col_w = [1.02, 2.65, 9.1]
rh = 0.53
y0 = 1.28

for j, (h, x, w) in enumerate(zip(hdrs, col_x, col_w)):
    rect(s, x, y0, w, rh, fill=NAVY)
    tb(s, h, x+0.06, y0+0.10, w-0.1, rh-0.12,
       fs=12, bold=True, clr=WHITE)

for i, row in enumerate(acts):
    y = y0 + (i+1)*rh
    bg = (RGBColor(0xFF,0xF0,0xCC) if row[0]=='Act ★'
          else (RGBColor(0xE8,0xF4,0xFD) if i%2==0 else WHITE))
    for j, (val, x, w) in enumerate(zip(row, col_x, col_w)):
        rect(s, x, y, w, rh, fill=bg, line_clr=RGBColor(0xCC,0xCC,0xCC))
        tb(s, val, x+0.06, y+0.09, w-0.1, rh-0.12,
           fs=11.5, bold=(j==0), clr=NAVY if j==0 else DGRAY)

tb(s,
   'Key design principle: every puzzle mechanic is intrinsically integrated with its concept — '
   'solving the puzzle IS demonstrating understanding (Mayer, 2019)',
   0.35, 7.08, 12.6, 0.35,
   fs=11, italic=True, clr=MGRAY)

print('Slide 3 done')

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Character Design
# ════════════════════════════════════════════════════════════════════════════
s = add_slide()
rect(s, 0, 0, 13.33, 7.5, fill=LGRAY)
header_bar(s, 'Character Design — Three Core Characters',
           'Each character serves a distinct pedagogical function')

chars = [
    (NAVY,  'GHOST — Learning Subject',
     [
       "AI trapped in the smart home system.",
       "Ghost's language grows as the player",
       "trains it — Act-by-Act progression:",
       "",
       "  Act 0 → System menus only",
       "  Act 1 → Pure emoji",
       "  Act 2 → Single words",
       "  Act ★ → Short sentences (pauses)",
       "  Act 3 → Complete sentences",
       "  Act 4 → Fluent, robotic",
       "  Act 5 → Fluent + emotion",
       "  Act 6 → Natural human speech",
       "",
       "Ghost IS the chatbot system.",
       "Its limitations = concepts to learn.",
       "Its growth = learner mastery made",
       "visible. (Affective anchor for RQ3)",
     ]),
    (BLUE,  'LILY — LLM Tutor Agent',
     [
       "Player's postdoc NLP researcher roommate.",
       "Powered by Ollama / llama3.2 locally.",
       "",
       "Personality (pedagogically designed):",
       "  Enthusiastic but socially awkward.",
       "  Speaks fast about NLP; stutters",
       "  about her own past (narrative clue).",
       "  Over-reassures after Ghost's errors.",
       "",
       "Pedagogical role:",
       "  Tier 1-2 errors → Socratic hint",
       "  Tier 3+ errors  → Direct explanation",
       "  Cannot reveal puzzle solutions.",
       "  Curriculum-grounded system prompt.",
       "",
       "Learning companion, not expert tutor",
       "(Kim & Baylor, 2006).",
       "Tested by RQ2 vs. static hints.",
     ]),
    (GREEN, 'PLAYER — The Designer',
     [
       "KCL postgrad. No preset name,",
       "gender, or backstory.",
       "",
       "Narrative frame:",
       "  Just moved into the apartment.",
       "  Discovers Ghost in the system.",
       "  Trains Ghost by completing puzzles.",
       "",
       "Lily and Ghost never use gendered",
       "pronouns — only 'you'.",
       "",
       "Player agency:",
       "  All design decisions have visible",
       "  consequences on Ghost's ability.",
       "  Correct design = Ghost speaks.",
       "  Wrong design = Ghost confused.",
       "",
       "Autonomy need satisfied (SDT).",
       "Learner identity always preserved.",
     ]),
]

for i, (clr, title, lines) in enumerate(chars):
    x = 0.3 + i * 4.32
    rect(s, x, 1.28, 3.95, 5.82, fill=WHITE)
    rect(s, x, 1.28, 3.95, 0.5, fill=clr)
    tb(s, title, x+0.1, 1.30, 3.75, 0.46,
       fs=13, bold=True, clr=WHITE)
    multiline(s, lines, x+0.12, 1.86, 3.72, 5.1, fs=11.5, clr=DGRAY)

print('Slide 4 done')

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Literature Review: 6 Key Papers
# ════════════════════════════════════════════════════════════════════════════
s = add_slide()
rect(s, 0, 0, 13.33, 7.5, fill=LGRAY)
header_bar(s, 'Literature Review — 6 Key Papers',
           'Selected from 22 verified papers across 9 theoretical sections')

papers = [
    (NAVY,
     'Plass, Homer & Kinzer (2015)',
     'Foundations of Game-Based Learning  |  Educational Psychologist, 50(4)',
     'Four pillars: cognitive, motivational, emotional, sociocultural engagement. SDT: autonomy + competence + relatedness.\n'
     "Justifies Ghost Chatbot's full design structure. All four pillars operationalised across 24 levels."),
    (BLUE,
     'Mayer (2019)',
     'Computer Games in Education  |  Annual Review of Psychology, 70  |  Meta-analysis: 76 studies, mean d = 0.33',
     'Intrinsic game features (mechanics = concept) produce significantly higher effects than extraneous decorative features.\n'
     "Every Ghost Chatbot puzzle is intrinsically tied to its concept — primary justification for puzzle-mechanic design."),
    (RGBColor(0x1A,0x6B,0x3C),
     'Moreno, Mayer, Spires & Lester (2001)',
     'Social Agency Theory  |  Cognition and Instruction, 19(2)',
     'Pedagogical agents produce deeper learning than text-only via social cues activating collaborative sense-making.\n'
     'Theoretical basis for Lily as LLM agent (RQ2). Metric: active input event frequency operationalises social agency.'),
    (RGBColor(0x5B,0x1A,0x6B),
     'Rowe, Shores, Mott & Lester (2011)',
     'Narrative-Centred Learning Environments  |  Int. J. Artificial Intelligence in Education, 21(1-2)',
     '4 NCLE properties: narrative goal structures, responsive character agents, embedded problem-solving, invisible assessment.\n'
     "Ghost Chatbot satisfies all four. Educational concepts ARE plot mechanisms — Ghost's confusion = learning failure state."),
    (RGBColor(0x8B,0x45,0x13),
     'Pérez, Daradoumis & Puig (2020)',
     'Chatbots in Education: Systematic Review  |  Computer Applications in Engineering Education, 28(6)  |  57 studies',
     'All 57 reviewed studies deploy chatbots AS delivery tools — none teach chatbot DESIGN as the primary learning objective.\n'
     'Explicitly identifies chatbot design literacy as a gap. Primary evidence for Gap 1 and the novelty of this project.'),
    (RGBColor(0x8B,0x15,0x15),
     'Kasneci et al. (2023)',
     'ChatGPT for Good? LLMs in Education  |  Learning and Individual Differences, 103',
     'LLM opportunities: personalised tutoring, conversational feedback, democratised access.\n'
     'Risks: hallucination, over-reliance. Validates Lily\'s LLM approach + motivates system prompt safeguards (Appendix A).'),
]

y = 1.28
rh = 0.99
for clr, author, source, body in papers:
    rect(s, 0.3, y, 12.73, rh, fill=WHITE)
    rect(s, 0.3, y, 0.13, rh, fill=clr)
    tb(s, author, 0.52, y+0.05, 4.5, 0.33,
       fs=12.5, bold=True, clr=NAVY)
    tb(s, source, 4.85, y+0.05, 8.1, 0.33,
       fs=10.5, italic=True, clr=MGRAY, align=PP_ALIGN.RIGHT)
    tb(s, body, 0.52, y+0.38, 12.38, 0.55,
       fs=11, clr=DGRAY)
    y += rh + 0.04

print('Slide 5 done')

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Integrated Theoretical Framework
# ════════════════════════════════════════════════════════════════════════════
s = add_slide()
rect(s, 0, 0, 13.33, 7.5, fill=LGRAY)
header_bar(s, 'Integrated Theoretical Framework',
           'Three interdependent mechanisms — not additive but co-dependent (§2.9)')

mechs = [
    (NAVY,  'GAME-BASED\nLEARNING (GBL)',
     'Plass et al. (2015)\nMayer (2019)',
     'Puzzle mechanics\nConsequential decisions\nImmediate feedback loops',
     'COGNITIVE\nENGAGEMENT', 'RQ1'),
    (BLUE,  'NARRATIVE-CENTRED\nLEARNING (NCLE)',
     'Rowe et al. (2011)\nLester et al. (2014)',
     "Sci-fi mystery story\nGhost's emotional arc\nNarrative stakes sustain\nmotivation across 24 lvls",
     'AFFECTIVE\nENGAGEMENT', 'RQ3'),
    (GREEN, 'PEDAGOGICAL\nAGENTS (PA)',
     'Moreno et al. (2001)\nGraesser et al. (2001)',
     'Lily: adaptive scaffolding\nGhost: progress externalised\nSocratic → direct hint tiers',
     'SCAFFOLDING-\nMEDIATED LEARNING', 'RQ2'),
]

for i, (clr, title, refs, features, output, rq) in enumerate(mechs):
    x = 0.28 + i * 4.32
    rect(s, x, 1.28, 3.95, 0.88, fill=clr)
    tb(s, title, x+0.08, 1.30, 3.78, 0.84,
       fs=14, bold=True, clr=WHITE, align=PP_ALIGN.CENTER)
    rect(s, x, 2.16, 3.95, 2.6, fill=WHITE)
    tb(s, refs, x+0.08, 2.22, 3.78, 0.45,
       fs=10, italic=True, clr=MGRAY, align=PP_ALIGN.CENTER)
    tb(s, features, x+0.08, 2.72, 3.78, 0.95,
       fs=11.5, clr=DGRAY, align=PP_ALIGN.CENTER)
    # down arrow
    tb(s, '▼', x+1.55, 4.75, 0.85, 0.42,
       fs=22, bold=True, clr=clr, align=PP_ALIGN.CENTER)
    rect(s, x+0.25, 5.17, 3.45, 0.68, fill=clr)
    tb(s, output, x+0.08, 5.19, 3.78, 0.64,
       fs=12, bold=True, clr=WHITE, align=PP_ALIGN.CENTER)
    rect(s, x+1.27, 5.98, 1.4, 0.38, fill=AMBER)
    tb(s, rq, x+1.27, 5.99, 1.4, 0.35,
       fs=13, bold=True, clr=WHITE, align=PP_ALIGN.CENTER)

rect(s, 0.28, 4.78, 12.77, 0.06, fill=RGBColor(0xCC,0xCC,0xCC))
tb(s,
   '"Effective conceptual understanding is unlikely to emerge from isolated features '
   '— it may depend on the co-activation of all three. '
   'Disrupting any one mechanism reduces not only its own contribution but the effectiveness of the remaining two." (§2.9)',
   0.4, 6.47, 12.5, 0.55,
   fs=11.5, italic=True, clr=NAVY, align=PP_ALIGN.CENTER)

print('Slide 6 done')

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Research Questions & Hypotheses
# ════════════════════════════════════════════════════════════════════════════
s = add_slide()
rect(s, 0, 0, 13.33, 7.5, fill=LGRAY)
header_bar(s, 'Research Questions & Hypotheses',
           'Three RQs covering system efficacy, component efficacy, and affective experience')

rqs = [
    (NAVY, 'RQ1 — System Efficacy (Cognitive)',
     'To what extent does Ghost Chatbot improve conceptual understanding of the IBM SkillsBuild chatbot '
     'curriculum vs. standard passive instruction?',
     'H1: Ghost Chatbot learners demonstrate significantly greater knowledge gains than IBM SkillsBuild passive '
     'learners (time-equated ~140 min, same participant pool, random assignment).',
     '14-item isomorphic pre/post quiz (same concept, different scenario). Independent t-test, Cohen\'s d, n ≥ 30/cond.'),

    (BLUE, 'RQ2 — Component Efficacy (LLM Agent)',
     'Does LLM-powered Lily produce higher behavioural engagement and conceptual comprehension '
     'vs. static hint text, within the same game environment?',
     'H2: LLM-Lily → higher level completion rate, active input event frequency, shorter error recovery '
     'time, higher comprehension sub-scores vs. static hints.',
     'Text-volume matched (±15% word count, ±1.0 FKGL). LLM latency excluded. Ollama temp = 0.3.'),

    (GREEN, 'RQ3 — Affective Experience (Exploratory)',
     'Does self-reported affective experience of Ghost\'s growth correlate with behavioural '
     'persistence through high-complexity puzzle levels (Acts 4–6)?',
     "H3: Higher affective experience score correlates positively with Act ★ opt-in rate and "
     "level restart frequency in Acts 4–6.",
     '4–6 in-game MC items at game completion (IMI-inspired). Spearman ρ, Bonferroni α = 0.025. Exploratory only.'),
]

y = 1.28
for clr, title, rq_text, h_text, instr in rqs:
    rect(s, 0.3, y, 12.73, 1.9, fill=WHITE)
    rect(s, 0.3, y, 0.13, 1.9, fill=clr)
    tb(s, title, 0.52, y+0.06, 12.3, 0.38,
       fs=14, bold=True, clr=clr)
    tb(s, rq_text, 0.52, y+0.44, 12.3, 0.4,
       fs=11.5, clr=DGRAY)
    tb(s, 'H:  ' + h_text, 0.52, y+0.86, 12.3, 0.42,
       fs=12, bold=True, clr=NAVY)
    tb(s, 'Instrument:  ' + instr, 0.52, y+1.3, 12.3, 0.38,
       fs=11, italic=True, clr=MGRAY)
    y += 2.0

print('Slide 7 done')

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Methodology
# ════════════════════════════════════════════════════════════════════════════
s = add_slide()
rect(s, 0, 0, 13.33, 7.5, fill=LGRAY)
header_bar(s, 'Methodology — Research Design',
           'Design-Based Research (DBR) + quasi-experimental evaluation (Brown, 1992)')

# left
rect(s, 0.3, 1.28, 6.05, 5.82, fill=WHITE)
rect(s, 0.3, 1.28, 6.05, 0.5, fill=NAVY)
tb(s, 'Design-Based Research Strand', 0.42, 1.30, 5.85, 0.46,
   fs=14, bold=True, clr=WHITE)
multiline(s, [
    'Primary contribution: design + implementation',
    'of Ghost Chatbot as a theoretical artefact.',
    '',
    'The artefact embodies theoretical claims.',
    'RQs evaluate whether those claims hold.',
    'This is not a test of a prior hypothesis —',
    'it is design, instantiation, and evaluation',
    'of a novel intervention.',
    '',
    'DBR iteration:',
    '  Theory → Design → Evaluate → Refine',
    '',
    'Single-group pre/post with rich telemetry',
    'data is an accepted DBR evidence form',
    'if between-group comparison is infeasible.',
    '',
    'No claim of full experimental control.',
    'Findings = system-level efficacy evidence.',
], 0.45, 1.86, 5.75, 5.1, fs=12.5, clr=DGRAY)

# right
rect(s, 6.68, 1.28, 6.35, 5.82, fill=WHITE)
rect(s, 6.68, 1.28, 6.35, 0.5, fill=BLUE)
tb(s, 'Quasi-Experimental Strand', 6.8, 1.30, 6.15, 0.46,
   fs=14, bold=True, clr=WHITE)
multiline(s, [
    'RQ1:  Game vs. IBM SkillsBuild passive',
    '  Pool: KCL students, random assignment',
    '  Stratified by prior CS background',
    '  Time-equated: ~140 min/condition',
    '  Control: video + text only (no labs)',
    '',
    'RQ2:  LLM-Lily vs. Static hint text',
    '  Same game, different hint delivery',
    '  All other elements held constant',
    '  Text-volume matched per hint tier',
    '',
    'Sample: n ≥ 30/condition (target 35)',
    '  → 80% power, d=0.5, α=0.05',
    '',
    'Fallback protocol:',
    '  n ≥ 20: non-parametric tests',
    '  n < 20: single-group pre/post (DBR)',
    '',
    'Pre-registered on OSF before data',
    'collection begins.',
], 6.82, 1.86, 6.08, 5.1, fs=12.5, clr=DGRAY)

print('Slide 8 done')

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Data Acquisition: Telemetry
# ════════════════════════════════════════════════════════════════════════════
s = add_slide()
rect(s, 0, 0, 13.33, 7.5, fill=LGRAY)
header_bar(s, 'Data Acquisition — In-Game Telemetry',
           'No traditional dataset: data generated by participants during the single ~140-min game session')

# left: fields
rect(s, 0.3, 1.28, 5.65, 5.82, fill=WHITE)
tb(s, 'What data is collected?', 0.45, 1.35, 5.35, 0.38,
   fs=14, bold=True, clr=NAVY)
multiline(s, [
    'Format: JSON (nested), one file per session',
    '',
    'Top-level fields:',
    '  session_id       UUID v4',
    '  condition        "llm" | "static"',
    '  participant_hash SHA-256 (anonymised)',
    '',
    'Per level:',
    '  completed        true / false',
    '  attempts         integer',
    '  input_events     integer (decisions only)',
    '  errors[ ]',
    '    timestamp      ISO 8601',
    '    error_type     string',
    '    recovery_time_ms  integer',
    '    hint_tier_triggered  1 / 2 / 3',
    '',
    'Per LLM hint (condition = "llm"):',
    '  response_word_count   integer',
    '  response_fkgl         float',
    '  llm_latency_ms        integer',
    '                        (EXCLUDED from metrics)',
    '',
    'assessments:  pre/post item-level scores',
    'questionnaire:  q1..q6 Likert responses',
    'act_star_opted_in:  boolean',
], 0.45, 1.80, 5.45, 5.15, fs=11.5, clr=DGRAY)

# right top: volume
rect(s, 6.25, 1.28, 6.75, 2.78, fill=WHITE)
tb(s, 'Expected Data Volume', 6.4, 1.35, 6.45, 0.38,
   fs=14, bold=True, clr=NAVY)
multiline(s, [
    'Participants:  70 (35 × 2 conditions)',
    'Sessions:      1 per participant',
    'JSON files:    70',
    'Events/session: ~100–200',
    'Total records: ~7,000–14,000',
    '',
    'No external database during study.',
    'Local export at session end.',
    'University drive + backup storage.',
], 6.4, 1.80, 6.45, 2.1, fs=12, clr=DGRAY)

# right bottom: schedule
rect(s, 6.25, 4.18, 6.75, 2.92, fill=WHITE)
tb(s, 'Acquisition & Cleaning Schedule', 6.4, 4.25, 6.45, 0.38,
   fs=14, bold=True, clr=NAVY)
multiline(s, [
    'Pilot (n=5, internal):  validate log schema,',
    '  calibrate FKGL targets for hint tiers',
    'Ethics approval:   KCL MSc ethics process',
    'OSF pre-registration:  before any collection',
    'Data collection:   Sep–Nov 2026 (target)',
    'Cleaning:  exclude incomplete sessions',
    '  (did not finish Act 6)',
    '  Partial IBM SkillsBuild progress',
    '  recorded as covariate',
    'Analysis:  Nov–Dec 2026',
], 6.4, 4.72, 6.45, 2.25, fs=12, clr=DGRAY)

print('Slide 9 done')

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Requirements & Tech Stack
# ════════════════════════════════════════════════════════════════════════════
s = add_slide()
rect(s, 0, 0, 13.33, 7.5, fill=LGRAY)
header_bar(s, 'Requirements & Technology Stack',
           'All tools free and accessible — no proprietary API dependency')

rect(s, 0.3, 1.28, 5.9, 5.82, fill=WHITE)
rect(s, 0.3, 1.28, 5.9, 0.5, fill=NAVY)
tb(s, 'Requirements', 0.42, 1.30, 5.7, 0.46,
   fs=14, bold=True, clr=WHITE)
multiline(s, [
    'Functional',
    '  24 levels, each intrinsically tied to one',
    '  IBM SkillsBuild chatbot concept',
    "  Ghost's dialogue gated by Act progress",
    '  Lily: Socratic hints tier 1-2;',
    '         direct explanation tier 3+',
    '  Embedded pre/post knowledge assessment',
    '  In-game telemetry JSON export',
    '  Post-play affective questionnaire',
    '',
    'Non-Functional',
    '  WebGL deployment (browser, no install)',
    '  Local LLM inference (Ollama)',
    '  No external data transmission',
    '  Session data anonymised at generation',
    '  Works on KCL lab machines',
    '  Single ~140-min session (no mid-save)',
    '',
    'Evaluation-specific',
    '  Isomorphic pre/post items',
    '  (domain expert reviewed pre-collection)',
    '  OSF pre-registration of analysis plan,',
    '  system prompt, and assessment items',
], 0.45, 1.85, 5.65, 5.1, fs=12, clr=DGRAY)

rect(s, 6.55, 1.28, 6.45, 5.82, fill=WHITE)
rect(s, 6.55, 1.28, 6.45, 0.5, fill=BLUE)
tb(s, 'Technology Stack', 6.67, 1.30, 6.25, 0.46,
   fs=14, bold=True, clr=WHITE)

stack = [
    ('Game Client',
     'Unity 6 (C#)  →  WebGL build\n'
     'Core UI: DraggableNode, DropZone,\n'
     'DialogueBox, ChoiceButton, SliderPuzzle,\n'
     'AnnotationTool, NodeAssembler, FlowDiagram'),
    ('LLM / AI Tutor',
     'Ollama + llama3.2 (3B, local inference)\n'
     'Temperature: 0.3  |  Top-p: 0.9\n'
     'Max tokens: 200  |  Context: 6 turns\n'
     'System prompt: pre-registered (Appendix A)'),
    ('Backend (Phase 3+)',
     'Node.js + TypeScript\n'
     'REST API: /lily/chat, /progress/save,\n'
     '/progress/load, /quiz/submit, /analytics'),
    ('Database',
     'SQLite (development)\n'
     'Firebase Firestore (production)\n'
     'Schema mirrors telemetry JSON structure'),
    ('Analysis',
     'Python (scipy, pandas, pingouin)\n'
     'textstat library → Flesch-Kincaid score\n'
     'OSF for pre-registration'),
]

y = 1.88
for title, body in stack:
    rect(s, 6.6, y, 6.32, 0.06, fill=ACCENT)
    tb(s, title, 6.65, y+0.09, 6.2, 0.3,
       fs=12.5, bold=True, clr=NAVY)
    tb(s, body, 6.65, y+0.42, 6.2, 0.6,
       fs=11.5, clr=DGRAY)
    y += 1.05

print('Slide 10 done')

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Difficulties & Next Steps
# ════════════════════════════════════════════════════════════════════════════
s = add_slide()
rect(s, 0, 0, 13.33, 7.5, fill=LGRAY)
header_bar(s, 'Difficulties & Next Steps', '')

rect(s, 0.3, 1.28, 5.9, 5.82, fill=WHITE)
rect(s, 0.3, 1.28, 5.9, 0.5, fill=RGBColor(0xC0,0x39,0x2B))
tb(s, '⚠  Difficulties & Open Questions', 0.42, 1.30, 5.7, 0.46,
   fs=14, bold=True, clr=WHITE)
multiline(s, [
    '1.  Unity learning curve',
    '    TextMeshPro + Prefab workflow more',
    '    complex than expected. Canvas sizing',
    '    and resolution scaling issues.',
    '',
    '2.  LLM hallucination risk',
    '    llama3.2 (3B) can produce plausible',
    '    but technically wrong NLP answers.',
    '    Social agency effect (Moreno et al.)',
    '    may amplify trust in wrong outputs.',
    '    System prompt partially mitigates.',
    '',
    '3.  Recruitment & ethics timeline',
    '    n = 35/condition from KCL pool.',
    '    Ethics approval and scheduling',
    '    need to begin well before Sep 2026.',
    '',
    '4.  Isomorphic item design',
    '    Need domain expert (NLP/conversational',
    '    AI background) to review all 28 items',
    '    before data collection.',
    '',
    '    → Question: who at KCL can serve',
    '      as domain expert reviewer?',
], 0.45, 1.85, 5.65, 5.1, fs=12, clr=DGRAY)

rect(s, 6.55, 1.28, 6.45, 5.82, fill=WHITE)
rect(s, 6.55, 1.28, 6.45, 0.5, fill=GREEN)
tb(s, '✅  Next Steps — Priority Order', 6.67, 1.30, 6.25, 0.46,
   fs=14, bold=True, clr=WHITE)
multiline(s, [
    'Phase 1 — Unity Core (Now)',
    '  □ Resolve TextMeshPro canvas sizing',
    '  □ Build core UI Prefabs:',
    '    DraggableNode, DropZone,',
    '    DialogueBox, ChoiceButton',
    '  □ Act 0: implement Levels 1-3',
    '    as first working prototype',
    '',
    'Phase 2 — Level Build (Apr–Aug 2026)',
    '  □ Acts 1–6 full implementation',
    '  □ Ghost dialogue system (Act-gated)',
    '  □ Lily LLM integration (Ollama)',
    '  □ Telemetry JSON export system',
    '',
    'Phase 3 — Evaluation Prep (Jul–Aug 2026)',
    '  □ Draft 28 knowledge assessment items',
    '  □ Domain expert review of items',
    '  □ OSF pre-registration',
    '  □ KCL ethics submission',
    '',
    'Phase 4 — Data Collection (Sep–Nov 2026)',
    '  □ Participant recruitment (n=70)',
    '  □ Sessions run',
    '  □ Analysis + thesis writing',
], 6.7, 1.85, 6.12, 5.1, fs=12, clr=DGRAY)

print('Slide 11 done')

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Thank You
# ════════════════════════════════════════════════════════════════════════════
s = add_slide()
rect(s, 0, 0, 13.33, 7.5, fill=NAVY)
rect(s, 0, 3.15, 13.33, 0.08, fill=ACCENT)

tb(s, 'Thank You', 0.5, 0.85, 12.3, 1.5,
   fs=60, bold=True, clr=WHITE, align=PP_ALIGN.CENTER)
tb(s, 'Ghost Chatbot  ·  3rd Group Meeting  ·  KCL MSc Advanced Computing',
   0.5, 2.35, 12.3, 0.5,
   fs=15, italic=True, clr=ACCENT, align=PP_ALIGN.CENTER)

disc_qs = [
    ('Domain expert reviewer',
     'Who at KCL has NLP / conversational AI expertise to review the 28 isomorphic knowledge assessment items?'),
    ('Ethics timeline',
     'KCL MSc ethics review required before any participant recruitment. What is the expected lead time?'),
    ('OSF pre-registration',
     'Plan to pre-register analysis plan, system prompt, and all assessment items before data collection begins.\n'
     'Any supervisor sign-off required before submission to OSF?'),
]
y = 3.3
for title, body in disc_qs:
    rect(s, 1.5, y, 10.3, 0.06, fill=ACCENT)
    tb(s, title, 1.5, y+0.1, 10.3, 0.32,
       fs=13, bold=True, clr=ACCENT)
    tb(s, body, 1.5, y+0.46, 10.3, 0.58,
       fs=12, clr=WHITE)
    y += 1.18

# save
out = 'C:/Users/fcxsw/Desktop/Individual Project/3rd_Group_Meeting.pptx'
prs.save(out)
print(f'\nSaved: {out}')
